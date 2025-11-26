import os
from functools import lru_cache
from pathlib import Path

from loguru import logger
from PIL.Image import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape as Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.table import Table, _Cell, _Row


class PPtExtractor:
    @lru_cache(maxsize=128)
    def generate_img_path(self, id: str, img_name: str) -> str:
        if not isinstance(id, str):
            raise ValueError("id must be a string")
        if not isinstance(img_name, str):
            raise ValueError("img_name must be a string")
        return f"media/{id}/{img_name}"

    # =========================================================================
    # Public Entry
    # =========================================================================

    def handle_shape(
        self,
        shape: Shape,
        content_list: list[dict[str, str]],
        media_dir: Path,
        img_map: dict[Path, str],
        id: str,
        skip_image: bool,
    ):
        """Dispatch shape to appropriate handler."""
        self._validate_args(shape, content_list, media_dir, img_map, id, skip_image)

        try:
            if self._is_text(shape):
                self._handle_text_shape(shape, content_list)
            elif self._is_picture(shape, skip_image):
                self._handle_picture_shape(shape, content_list, media_dir, img_map, id)
            elif self._is_table(shape):
                self._handle_table_shape(shape, content_list)
            elif self._is_group(shape):
                self._handle_group_shape(
                    shape, content_list, media_dir, img_map, id, skip_image
                )
            else:
                logger.info(f"Unknown shape type: {shape.shape_type}, {type(shape)}")

        except (PermissionError, OSError) as e:
            logger.error(f"File error handling shape: {e}")
        except Exception as e:
            logger.error(f"Unexpected error handling shape: {e}")

    # =========================================================================
    # Type Detection
    # =========================================================================

    def _is_text(self, shape: Shape) -> bool:
        return hasattr(shape, "has_text_frame") and shape.has_text_frame

    def _is_picture(self, shape: Shape, skip_image: bool) -> bool:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not skip_image

    def _is_table(self, shape: Shape) -> bool:
        return shape.shape_type == MSO_SHAPE_TYPE.TABLE

    def _is_group(self, shape: Shape) -> bool:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP

    # =========================================================================
    # Handlers
    # =========================================================================

    # --- Text ---------------------------------------------------------------

    def _handle_text_shape(self, shape: Shape, content_list: list):
        for paragraph in shape.text_frame.paragraphs:
            content_list.append({"type": "text", "data": paragraph.text + "\n"})

    # --- Picture ------------------------------------------------------------

    def _handle_picture_shape(
        self,
        shape: Picture,
        content_list: list,
        media_dir: Path,
        img_map: dict[Path, str],
        id: str,
    ):
        image: Image = shape.image
        img_path = self._prepare_image_path(media_dir, img_map, image.ext)
        img_s3_path = self.generate_img_path(id, img_path.name)

        # map local→remote
        img_map[img_path] = img_s3_path

        # record content
        content_list.append({"type": "image", "data": img_s3_path})

        # write file
        with open(img_path, "wb") as f:
            f.write(image.blob)

    def _prepare_image_path(self, media_dir: Path, img_map: dict, ext: str) -> Path:
        if not media_dir.exists():
            media_dir.mkdir(parents=True, exist_ok=True)

        if not os.access(media_dir, os.W_OK):
            raise PermissionError(f"Cannot write to directory: {media_dir}")

        filename = f"pic-{len(img_map)}.{ext}"
        return media_dir / filename

    # --- Table --------------------------------------------------------------

    def _handle_table_shape(self, shape: GraphicFrame, content_list: list):
        table: Table = shape.table
        md = self._table_to_markdown(table)
        content_list.append({"type": "md", "data": md})

    def _table_to_markdown(self, table: Table) -> str:
        md = "\n"
        for row_idx, row in enumerate(table.rows):
            md += "|"
            if row_idx == 1:  # header separator
                for _ in row.cells:
                    md += "---|"
                md += "\n|"

            for cell in row.cells:
                text = cell.text.replace("\r", " ").replace("\n", " ")
                md += f" {text} |"

            md += "\n"
        md += "\n"
        return md

    # --- Group --------------------------------------------------------------

    def _handle_group_shape(
        self,
        shape: GroupShape,
        content_list: list,
        media_dir: Path,
        img_map: dict,
        id: str,
        skip_image: bool,
    ):
        for sub_shape in shape.shapes:
            self.handle_shape(
                sub_shape,
                content_list,
                media_dir,
                img_map,
                id,
                skip_image,
            )

    # ======================================================================
    # Public API
    # ======================================================================

    def extract(
        self,
        presentation_source: Path,
        id: str,
        dir: Path,
        media_dir: Path,
        skip_image: bool,
    ):
        """Extract content from PPTX."""
        self._validate_args(presentation_source, id, dir, media_dir, skip_image)

        return self._extract_pages(presentation_source, id, media_dir, skip_image)

    # ======================================================================
    # Validation helpers
    # ======================================================================

    def _validate_args(
        self,
        presentation_source: Path,
        id: str,
        dir: Path,
        media_dir: Path,
        skip_image: bool,
    ):
        validators = [
            (presentation_source, Path, "presentation_source must be a Path object"),
            (id, str, "id must be a string"),
            (dir, Path, "dir must be a Path object"),
            (media_dir, Path, "media_dir must be a Path object"),
            (skip_image, bool, "skip_image must be a boolean"),
        ]

        for value, type_, msg in validators:
            if not isinstance(value, type_):
                raise ValueError(msg)

    # ======================================================================
    # Extract core
    # ======================================================================

    def _extract_pages(
        self,
        presentation_source: Path,
        id: str,
        media_dir: Path,
        skip_image: bool,
    ):
        pages = []
        img_map = {}

        try:
            presentation = Presentation(presentation_source)
            for page_no, slide in enumerate(presentation.slides):
                page = self._extract_page(
                    slide, page_no, id, media_dir, img_map, skip_image
                )
                pages.append(page)

        except (FileNotFoundError, PermissionError, OSError) as e:
            logger.error(f"File error: {e}")
        except Exception as e:
            logger.error(f"Unexpected error extracting presentation: {e}")

        return pages

    # ======================================================================
    # Extract a single page
    # ======================================================================

    def _extract_page(
        self,
        slide: Slide,
        page_no: int,
        id: str,
        media_dir: Path,
        img_map: dict,
        skip_image: bool,
    ):
        page = {"page_no": page_no, "content_list": []}

        for shape in slide.shapes:
            self.handle_shape(
                shape,
                page["content_list"],
                media_dir,
                img_map,
                id,
                skip_image,
            )

        return page

    def run(self, id: str, file_path: Path, skip_image: bool = False):
        if not isinstance(id, str):
            raise ValueError("id must be a string")
        if not isinstance(file_path, Path):
            raise ValueError("file_path must be a Path object")
        if not isinstance(skip_image, bool):
            raise ValueError("skip_image must be a boolean")

        media_dir = Path("media").resolve()
        return self.extract(file_path, id, Path("."), media_dir, skip_image)
